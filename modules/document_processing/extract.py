"""
Extract structured text from downloaded GOJEP tender documents.

ONE extraction module for both tenders and email workflows.

Format routing:
  PDF, DOCX        -> Lightning AI Studio (Docling, GPU)
  XLSX, XLS        -> pandas (sheet-level extraction, local)
  PPTX, PPT        -> python-pptx (slide-level extraction, local)
  XML, TXT, others -> raw text read (local)

Folder structure per tender:
  <tender_id>/
      tender_data/
          document_downloads/    <- nested folder structure, all unzipped files
          extracted_documents/   <- flat copy of files (only changed since last sync)
          json_documents/        <- extracted JSON output + .manifest.json
      email_updates/
          clarifications/
              document_downloads/
              extracted_documents/
              json_documents/
          new_documents/
              document_downloads/
              extracted_documents/
              json_documents/

For PDF/DOCX the flow is:
  1. Upload file to Supabase Storage (tender-documents bucket)
  2. Upsert a pending row in document_extractions
  3. Dispatch Lightning Studio job (runs Docling on all pending files)
  4. Poll document_extractions until status = completed / failed
  5. Save extracted markdown to json_documents/<filename>.json
"""

from __future__ import annotations

import json
import logging
import os
import shutil
import time
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List

import pandas as pd
from pptx import Presentation

from config import settings as config
from modules.shared.document_sync import (
    build_file_manifest,
    load_extraction_manifest,
    save_extraction_manifest,
    update_file_in_manifest,
    compare_manifests,
    TENDER_DATA,
    DOCUMENT_DOWNLOADS,
    EXTRACTED_DOCUMENTS,
    JSON_DOCUMENTS,
)

logger = logging.getLogger(__name__)


class SafeEncoder(json.JSONEncoder):
    def default(self, obj):
        try:
            return super().default(obj)
        except TypeError:
            return str(obj)


DOCLING_FORMATS = {".pdf", ".docx"}
XLSX_FORMATS = {".xlsx", ".xls"}
PPTX_FORMATS = {".pptx", ".ppt"}
SKIP_EXTENSIONS = {".zip", ".json", ".jpg", ".png", ".jpeg", ".log", ".gif", ".bmp", ".tiff"}

_BUCKET = "tender-documents"
_TABLE = "document_extractions"

_CONTENT_TYPES = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xls": "application/vnd.ms-excel",
}


def _sanitize_storage_path(path: str) -> str:
    """Sanitize each component of a relative storage path for Supabase Storage."""
    def _sanitize(part: str) -> str:
        safe = ""
        for ch in part:
            if ch.isascii() and (ch.isalnum() or ch in "-_.() "):
                safe += ch
            else:
                safe += "_"
        return safe

    return "/".join(_sanitize(p) for p in path.split("/"))


def _get_ext(file_path: str) -> str:
    """Return lowercase extension with leading dot (e.g. '.pdf')."""
    _, ext = os.path.splitext(file_path.lower())
    return ext


# ── Lightning Studio client ──────────────────────────────────────────────────

def _get_lightning_studio():
    """Return a connected Lightning Studio instance."""
    from lightning_sdk import Studio

    studio = Studio(
        name="docling-gpu-extraction-devbox",
        teamspace="document-information-extraction-project",
        user="klearrshipping",
    )
    return studio


def dispatch_to_lightning():
    """
    Start the Lightning Studio and submit the extraction job.
    Does NOT stop the studio — the caller is responsible for stopping it
    after confirming results are in Supabase.
    """
    studio = _get_lightning_studio()

    print("Connecting to Lightning Studio...")

    # Poll until studio leaves transient states (Pending, Stopping)
    # Calling start() during Pending raises a RuntimeError
    max_wait_secs = 120
    poll_interval  = 5
    waited = 0
    while True:
        status = str(getattr(studio, "status", "") or "").lower()
        if status not in ("pending", "stopping"):
            break
        if waited >= max_wait_secs:
            raise RuntimeError(
                f"Lightning Studio stuck in '{status}' for {max_wait_secs}s — "
                "cancel or terminate the job in the Lightning UI before retrying."
            )
        print(f"  Studio is {status}, waiting {poll_interval}s...")
        time.sleep(poll_interval)
        waited += poll_interval

    # Now status is stable — act accordingly
    if status == "running":
        print("Studio already running.")
    elif status in ("stopped", ""):
        print("Studio is stopped. Starting...")
        try:
            studio.start()
            print("Studio started.")
        except Exception as e:
            # Check both the exception message and any HTTP response body
            err_str = str(e)
            err_body = str(getattr(e, "body", "") or "")
            if "already has instances running" in err_str or "already has instances running" in err_body:
                # A stale instance (e.g. cpu-4) is still alive — stop it and retry
                print(f"  Stale instance detected — stopping before retry...")
                try:
                    studio.stop()
                except Exception:
                    pass
                # Wait for it to fully stop
                for _ in range(24):
                    s = str(getattr(studio, "status", "") or "").lower()
                    if s == "stopped":
                        break
                    print(f"  Waiting for studio to stop (status: {s})...")
                    time.sleep(5)
                # Wait for it to fully stop before retrying
                for _ in range(24):
                    s = str(getattr(studio, "status", "") or "").lower()
                    if s not in ("pending", "stopping"):
                        break
                    print(f"  Waiting for studio to reach stopped state (status: {s})...")
                    time.sleep(5)
                print("  Retrying start...")
                try:
                    studio.start()
                    print("Studio started.")
                except Exception as retry_err:
                    raise RuntimeError(
                        f"Lightning Studio failed to start after clearing stale instance: {retry_err}"
                    ) from retry_err
            else:
                raise
    else:
        raise RuntimeError(f"Lightning Studio in unexpected state: '{status}'. Manual intervention required.")

    # Wait for the studio to fully initialise before submitting a job.
    # After start(), the studio may still be in a "setting things up" phase
    # even though status shows "running". Poll until studio.run() succeeds.
    script_path = "/teamspace/studios/this_studio/main.py"
    print(f"Running extraction script on Lightning: {script_path}")
    submitted = False
    for attempt in range(12):  # retry up to ~60s
        try:
            studio.run(f"python {script_path}")
            print("Lightning job submitted.")
            submitted = True
            break
        except Exception as e:
            err = str(e)
            err_body = str(getattr(e, "body", "") or "")
            retriable = (
                "still setting things up" in err
                or "progress bar" in err
                or "response: 500" in err
                or "response: 500" in err_body
            )
            if retriable:
                print(f"  Studio not ready, retrying in 10s... (attempt {attempt + 1}/12)")
                time.sleep(10)
            else:
                print(f"Lightning job submission failed: {e}")
                try:
                    studio.stop()
                except Exception:
                    pass
                return None
    if not submitted:
        print("Lightning studio failed to initialise after 2 minutes — skipping.")
        try:
            studio.stop()
        except Exception:
            pass
        return None


def stop_lightning_studio(studio) -> None:
    """Stop a Lightning Studio instance."""
    try:
        studio.stop()
        print("Studio stopped.")
    except Exception as e:
        print(f"Could not stop studio: {e}")


# ── Lightning Studio extraction (PDF / DOCX) ──────────────────────────────────

def extract_via_lightning_flat(
    local_paths: List[str],
    json_docs_dir: str,
    poll_interval: int = 30,
    timeout: int = 1800,
) -> int:
    """
    Upload local PDF/DOCX files to Supabase Storage, dispatch Lightning Studio
    job, poll until complete, then save JSON results to json_docs_dir.
    All files belong to the same tender folder, so json_docs_dir is a single path.
    """
    import requests as rq

    if not local_paths:
        return 0

    supabase_url = config.SUPABASE_URL.rstrip("/")
    key = config.SUPABASE_SECRET_KEY

    storage_headers = {
        "apikey": key,
        "Authorization": f"Bearer {key}",
    }
    rest_headers = {
        "apikey": key,
        "Authorization": f"Bearer {key}",
        "Content-Type": "application/json",
        "Prefer": "resolution=merge-duplicates,return=minimal",
    }

    # json_docs_dir = .../documents/<tender_id>/tender_data/json_documents
    # We need the competition_unique_id folder name, two levels up from json_documents
    tender_id = Path(json_docs_dir).parent.parent.name

    pending: list[tuple[str, str]] = []

    for local_path in local_paths:
        raw_filename = Path(local_path).name
        sanitized_fname = _sanitize_storage_path(raw_filename)
        storage_path = f"{tender_id}/{sanitized_fname}"
        suffix = Path(local_path).suffix.lower()
        content_type = _CONTENT_TYPES.get(suffix, "application/octet-stream")

        print(f"    Uploading: {raw_filename}", flush=True)
        try:
            with open(local_path, "rb") as f:
                resp = rq.post(
                    f"{supabase_url}/storage/v1/object/{_BUCKET}/{storage_path}",
                    headers={
                        **storage_headers,
                        "Content-Type": content_type,
                        "x-upsert": "true",
                    },
                    data=f,
                    timeout=60,
                )
            if resp.status_code not in (200, 201):
                logger.error("Storage upload failed (%s): %s — %s",
                             resp.status_code, storage_path, resp.text[:200])
                continue

            rq.post(
                f"{supabase_url}/rest/v1/{_TABLE}",
                headers=rest_headers,
                params={"on_conflict": "tender_resource_id,filename"},
                json={
                    "tender_resource_id": tender_id,
                    "filename": sanitized_fname,
                    "storage_path": storage_path,
                    "extraction_status": "pending",
                },
                timeout=30,
            )
            pending.append((local_path, sanitized_fname))

        except Exception as e:
            logger.error("Upload error for %s: %s", local_path, e)

    if not pending:
        logger.warning("No files uploaded — Lightning extraction skipped.")
        return 0

    print(f"    Dispatching Lightning Studio for {len(pending)} file(s)...", flush=True)
    studio = dispatch_to_lightning()

    print("    Polling for extraction results...", flush=True)
    poll_headers = {
        "apikey": key,
        "Authorization": f"Bearer {key}",
        "Content-Type": "application/json",
    }
    deadline = time.time() + timeout
    completed: dict[str, str | None] = {}

    try:
        while time.time() < deadline:
            for local_path, sanitized_fname in pending:
                if local_path in completed:
                    continue
                try:
                    resp = rq.get(
                        f"{supabase_url}/rest/v1/{_TABLE}",
                        headers=poll_headers,
                        params={
                            "tender_resource_id": f"eq.{tender_id}",
                            "filename": f"eq.{sanitized_fname}",
                            "select": "extraction_status,extracted_text",
                        },
                        timeout=30,
                    )
                    if resp.ok:
                        rows = resp.json()
                        if rows:
                            status = rows[0].get("extraction_status")
                            if status == "completed":
                                completed[local_path] = rows[0].get("extracted_text", "")
                            elif status == "failed":
                                completed[local_path] = None
                                logger.error("Lightning extraction failed for %s/%s",
                                             tender_id, sanitized_fname)
                except Exception as e:
                    logger.warning("Poll error for %s: %s", sanitized_fname, e)

            if len(completed) == len(pending):
                break

            print(f"    Waiting... ({len(completed)}/{len(pending)} complete)", flush=True)
            time.sleep(poll_interval)
        else:
            logger.warning("Lightning extraction timed out after %ss (%d/%d complete)",
                           timeout, len(completed), len(pending))
    finally:
        if studio:
            print("    Stopping Lightning Studio...", flush=True)
            stop_lightning_studio(studio)

    success_count = 0
    for local_path, extracted_text in completed.items():
        if extracted_text is None:
            continue

        file_name = os.path.basename(local_path)
        output_json_path = os.path.join(json_docs_dir, f"{file_name}.json")

        ext = _get_ext(local_path)
        output_data = {
            "source_file": file_name,
            "extension": ext.lstrip("."),
            "extraction_timestamp": datetime.utcnow().isoformat() + "Z",
            "content": {
                "markdown": extracted_text,
                "source": "lightning",
            },
        }

        try:
            with open(output_json_path, "w", encoding="utf-8") as f:
                json.dump(output_data, f, ensure_ascii=False, indent=2, cls=SafeEncoder)
            print(f"    Saved: {file_name}.json", flush=True)
            success_count += 1
        except Exception as e:
            logger.error("Failed to save JSON for %s: %s", file_name, e)

    return success_count


# ── Local extractors ─────────────────────────────────────────────────────────

def extract_xlsx(file_path: str) -> Dict[str, Any]:
    """Load all sheets, normalise NaNs to empty strings."""
    try:
        dfs = pd.read_excel(file_path, sheet_name=None)
        return {
            sheet: df.fillna("").astype(str).to_dict(orient="records")
            for sheet, df in dfs.items()
        }
    except Exception as e:
        logger.error(f"XLSX extraction failed for {file_path}: {e}")
        return {"error": str(e)}


def extract_pptx(file_path: str) -> Dict[str, Any]:
    """Extract slide text and tables from a PowerPoint presentation."""
    try:
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            tables = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
                if shape.has_table:
                    tables.append([
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ])
            slides.append({"slide_number": i, "texts": texts, "tables": tables})
        return {"slides": slides}
    except Exception as e:
        logger.error(f"PPTX extraction failed for {file_path}: {e}")
        return {"error": str(e)}


def extract_text(file_path: str) -> Dict[str, Any]:
    """Read file as plain text."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return {"raw_text": f.read()}
    except Exception as e:
        logger.error(f"Text extraction failed for {file_path}: {e}")
        return {"error": str(e)}


def _process_single_file(file_path: str, json_docs_dir: str) -> bool:
    """Extract a single non-PDF/DOCX file and save JSON to json_docs_dir."""
    ext = _get_ext(file_path)
    file_name = os.path.basename(file_path)
    output_json_path = os.path.join(json_docs_dir, f"{file_name}.json")

    if ext in XLSX_FORMATS:
        content = extract_xlsx(file_path)
    elif ext in PPTX_FORMATS:
        content = extract_pptx(file_path)
    else:
        content = extract_text(file_path)

    if content is None:
        return False

    if list(content.keys()) == ["error"]:
        logger.warning(f"Extraction failed for {file_name}: {content['error']}")
        return False

    output_data = {
        "source_file": file_name,
        "extension": ext.lstrip("."),
        "extraction_timestamp": datetime.utcnow().isoformat() + "Z",
        "content": content,
    }

    try:
        with open(output_json_path, "w", encoding="utf-8") as f:
            json.dump(output_data, f, ensure_ascii=False, indent=2, cls=SafeEncoder)
        return True
    except Exception as e:
        logger.error(f"Failed to save JSON for {file_name}: {e}")
        return False


# ── Recursive ZIP extraction ───────────────────────────────────────────────────

def extract_nested_zips(target_folder: str, source_subfolder: str) -> None:
    """
    Recursively extract any .zip files found in
    <target_folder>/<source_subfolder>/document_downloads/.
    Zips are extracted in-place and deleted.
    """
    downloads_dir = os.path.join(target_folder, source_subfolder, DOCUMENT_DOWNLOADS)
    if not os.path.exists(downloads_dir):
        return

    changed = True
    while changed:
        changed = False
        for root, dirs, files in os.walk(downloads_dir):
            dirs[:] = [d for d in dirs if d not in (EXTRACTED_DOCUMENTS, JSON_DOCUMENTS)]
            for zip_name in files:
                if not zip_name.lower().endswith(".zip"):
                    continue
                zip_path = os.path.join(root, zip_name)
                extract_dir = os.path.splitext(zip_path)[0]
                try:
                    os.makedirs(extract_dir, exist_ok=True)
                    with zipfile.ZipFile(zip_path, "r") as zf:
                        zf.extractall(extract_dir)
                    os.remove(zip_path)
                    logger.info("Extracted and removed nested zip: %s", zip_path)
                    changed = True
                except zipfile.BadZipFile:
                    logger.warning("Bad or corrupted zip skipped: %s", zip_path)
                    os.remove(zip_path)
                except Exception as e:
                    logger.warning("Failed to extract nested zip %s: %s", zip_path, e)


# ── Sync and extract ─────────────────────────────────────────────────────────

def sync_and_extract(
    target_folder: str,
    source_subfolder: str,
) -> dict:
    """
    Extract only new/changed files from <source_subfolder>/document_downloads/
    directly to <source_subfolder>/json_documents/.

    document_downloads/ is already flat (ZIP extraction flattens at download
    time), so no intermediate extracted_documents/ copy is needed.

    Returns:
        {newly_processed, skipped, errors}
    """
    downloads_dir = os.path.join(target_folder, source_subfolder, DOCUMENT_DOWNLOADS)
    json_dir      = os.path.join(target_folder, source_subfolder, JSON_DOCUMENTS)

    if not os.path.exists(downloads_dir):
        return {"newly_processed": 0, "skipped": 0, "errors": 0}

    if os.path.isfile(json_dir):
        os.remove(json_dir)
    os.makedirs(json_dir, exist_ok=True)

    # Build manifest of all files in document_downloads/ (recursive for any
    # legacy tenders that still have a nested structure from old downloads)
    original_manifest   = build_file_manifest(downloads_dir, relative_to=downloads_dir)
    extraction_manifest = load_extraction_manifest(json_dir)
    existing_files      = extraction_manifest.get("files", {})

    comparison    = compare_manifests(existing_files, original_manifest)
    files_to_sync = [fname for fname, status in comparison.items() if status in ("new", "changed")]

    if not files_to_sync:
        logger.info(f"  No files need syncing in {source_subfolder}")
        return {"newly_processed": 0, "skipped": 0, "errors": 0}

    # Log new/changed files
    for fname in files_to_sync:
        logger.info(f"  New/changed: {fname}")

    total     = len(files_to_sync)
    processed = 0
    skipped   = 0
    errors    = 0

    print(f"  Found {total} files to extract in {source_subfolder}.", flush=True)

    lightning_batch: list[tuple[str, str]] = []

    for i, fname in enumerate(files_to_sync, start=1):
        # fname may be a relative path for legacy nested tenders; basename for new flat ones
        file_name = os.path.basename(fname)
        file_path = os.path.join(downloads_dir, fname)
        print(f"  [{i}/{total}] {file_name}", flush=True)
        ext = _get_ext(file_path)

        if ext in SKIP_EXTENSIONS:
            skipped += 1
            continue

        output_json_path = os.path.join(json_dir, f"{file_name}.json")
        if os.path.exists(output_json_path):
            skipped += 1
            continue

        if ext in DOCLING_FORMATS:
            lightning_batch.append((file_path, json_dir))
        else:
            try:
                if _process_single_file(file_path, json_dir):
                    processed += 1
                else:
                    skipped += 1
            except Exception as e:
                print(f"    -> ERROR: {e}", flush=True)
                logger.error(f"Failed to process {file_path}: {e}")
                errors += 1

    if lightning_batch:
        print(f"\n  Dispatching {len(lightning_batch)} PDF/DOCX file(s) to Lightning...", flush=True)
        paths = [p for p, _ in lightning_batch]
        extracted = extract_via_lightning_flat(paths, json_docs_dir=json_dir)
        processed += extracted
        errors += len(lightning_batch) - extracted

    # Update manifest for all processed files
    for fname in files_to_sync:
        src_path = os.path.join(downloads_dir, fname)
        if os.path.exists(src_path):
            stat    = os.stat(src_path)
            hash_md5 = original_manifest.get(fname, {}).get("hash_md5", "")
            update_file_in_manifest(
                extraction_manifest,
                fname,
                stat.st_size,
                hash_md5,
                f"{source_subfolder}/{DOCUMENT_DOWNLOADS}/{fname}",
            )

    save_extraction_manifest(json_dir, extraction_manifest)

    logger.info(f"  Extraction complete for {source_subfolder}: processed={processed}, skipped={skipped}, errors={errors}")
    return {"newly_processed": processed, "skipped": skipped, "errors": errors}


# ── Entry points ──────────────────────────────────────────────────────────────

def run_tender_extraction(tender_folder: str) -> dict:
    """
    Extract documents for a single tender from tender_data/ folder.
    document_downloads/ is already flat — ZIP extraction flattens at download time.
    """
    return sync_and_extract(tender_folder, TENDER_DATA)


def run_email_extraction(tender_folder: str, email_subfolder: str) -> dict:
    """
    Extract documents for a single email update subfolder.
    document_downloads/ is already flat — ZIP extraction flattens at download time.
    """
    return sync_and_extract(tender_folder, email_subfolder)


def run_document_extraction() -> dict:
    """
    Process ALL tender folders. For each tender:
      - Extract tender_data documents
      - Extract email_updates/clarifications documents
      - Extract email_updates/new_documents documents

    This is the standalone CLI entry point.
    """
    from modules.shared.document_sync import EMAIL_UPDATES, CLARIFICATIONS, NEW_DOCUMENTS

    docs_dir = os.path.join(config.TENDERS_OUTPUT_DIRECTORY, "documents")
    if not os.path.exists(docs_dir):
        logger.warning(f"Documents directory not found at {docs_dir}")
        return {"total_files_scanned": 0, "newly_processed": 0, "skipped": 0, "errors": 0}

    tender_folders = [
        d for d in os.listdir(docs_dir)
        if os.path.isdir(os.path.join(docs_dir, d)) and not d.startswith("_")
    ]

    total_processed = 0
    total_skipped = 0
    total_errors = 0

    for tender_folder in tender_folders:
        tender_path = os.path.join(docs_dir, tender_folder)
        logger.info(f"\nProcessing tender: {tender_folder}")

        # Extract tender documents
        result = run_tender_extraction(tender_path)
        total_processed += result["newly_processed"]
        total_skipped += result["skipped"]
        total_errors += result["errors"]

        # Extract email update documents
        for email_subfolder in [EMAIL_UPDATES]:
            email_path = os.path.join(tender_path, email_subfolder)
            if not os.path.exists(email_path):
                continue

            for sub in [CLARIFICATIONS, NEW_DOCUMENTS]:
                sub_path = os.path.join(email_path, sub)
                if os.path.exists(sub_path):
                    result = run_email_extraction(tender_path, os.path.join(email_subfolder, sub))
                    total_processed += result["newly_processed"]
                    total_skipped += result["skipped"]
                    total_errors += result["errors"]

    return {
        "total_files_scanned": total_processed + total_skipped,
        "newly_processed": total_processed,
        "skipped": total_skipped,
        "errors": total_errors,
    }


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    run_document_extraction()