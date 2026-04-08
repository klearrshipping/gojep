"""
Extract structured text from downloaded GOJEP tender documents.

Format routing:
  PDF, DOCX        -> Docling via WSL venv (structure-aware, chunked output)
  XLSX, XLS        -> pandas (sheet-level extraction)
  PPTX, PPT        -> python-pptx (slide-level extraction)
  XML, TXT, others -> raw text read

Output: <tender_folder>/extracted_docs/<filename>.json
"""

import json
import logging
import os
import subprocess
import tempfile
from datetime import datetime
from typing import Any, Dict

import pandas as pd
from pptx import Presentation

from config import settings as config

logger = logging.getLogger(__name__)


class SafeEncoder(json.JSONEncoder):
    def default(self, obj):
        try:
            return super().default(obj)
        except TypeError:
            return str(obj)


DOCLING_FORMATS = {".pdf", ".docx"}
XLSX_FORMATS = {".xlsx", ".xls"}
PPTX_FORMATS = {".pptx", ".ppt"}
SKIP_EXTENSIONS = {".zip", ".json", ".jpg", ".png", ".jpeg", ".log", ".gif", ".bmp", ".tiff"}

# WSL Python path — loaded from config so it can be overridden via env var
WSL_PYTHON = config.WSL_PYTHON


def _get_ext(file_path: str) -> str:
    """Return lowercase extension with leading dot (e.g. '.pdf')."""
    _, ext = os.path.splitext(file_path.lower())
    return ext


def _windows_path_to_wsl(path: str) -> str:
    """Convert a Windows absolute path to its WSL /mnt/... equivalent."""
    path = path.replace("\\", "/")
    if len(path) >= 2 and path[1] == ":":
        drive = path[0].lower()
        path = f"/mnt/{drive}/{path[3:]}"
    return path


# ── Docling via WSL persistent worker (PDF + DOCX) -------------------------
#
# A single WSL Python process starts once, loads the docling models, then
# reads file paths from stdin and writes JSON results to stdout — one line
# per file. This avoids the 30-60s model re-initialisation cost per file.

_DOCLING_DAEMON = """
import sys, json, warnings, re
warnings.filterwarnings("ignore")

from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions

opts = PdfPipelineOptions()
opts.do_ocr = False
opts.do_table_structure = False
opts.do_picture_classification = False
opts.do_picture_description = False

converter = DocumentConverter(
    format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=opts)}
)

# Signal readiness
print("READY", flush=True)

for line in sys.stdin:
    file_path = line.strip()
    if not file_path:
        continue
    try:
        result = converter.convert(file_path)
        doc = result.document
        markdown = doc.export_to_markdown()

        sections = re.split(r'(?=^## )', markdown, flags=re.MULTILINE)
        chunks = []
        current_headings = []
        for section in sections:
            if not section.strip():
                continue
            lines = section.strip().splitlines()
            heading = lines[0].lstrip("#").strip() if lines[0].startswith("#") else None
            if heading:
                current_headings = [heading]
            chunks.append({
                "text": section.strip(),
                "headings": list(current_headings),
                "page": None,
            })

        output = {"ok": True, "markdown": markdown, "chunks": chunks}
    except Exception as e:
        output = {"ok": False, "error": str(e)}

    print(json.dumps(output, ensure_ascii=False), flush=True)
"""


class _DoclingWorker:
    """Persistent WSL process that keeps docling models loaded between files."""

    def __init__(self):
        self._proc = None

    def _ensure_started(self):
        if self._proc and self._proc.poll() is None:
            return
        print("  [docling] Starting WSL worker (loading models)...", flush=True)

        # Write daemon script to a temp file so we avoid Windows command-line
        # length limits and quoting issues when passing a long -c argument.
        self._daemon_file = tempfile.NamedTemporaryFile(
            mode="w", suffix=".py", delete=False, encoding="utf-8"
        )
        self._daemon_file.write(_DOCLING_DAEMON)
        self._daemon_file.close()
        wsl_script = _windows_path_to_wsl(self._daemon_file.name)

        self._proc = subprocess.Popen(
            ["wsl", "-d", "Ubuntu", "--", WSL_PYTHON, wsl_script],
            stdin=subprocess.PIPE,
            stdout=subprocess.PIPE,
            stderr=subprocess.DEVNULL,   # discard stderr — piping it causes deadlock
            text=True,
            encoding="utf-8",
        )
        # Wait for READY signal (up to 120s for model loading)
        import threading
        ready_line = []
        def _read_ready():
            ready_line.append(self._proc.stdout.readline())
        t = threading.Thread(target=_read_ready, daemon=True)
        t.start()
        t.join(timeout=120)
        if t.is_alive():
            self._proc.kill()
            raise RuntimeError("Docling worker timed out waiting for READY signal")

        ready = ready_line[0].strip().replace("\x00", "") if ready_line else ""
        if ready != "READY":
            raise RuntimeError(
                f"Docling worker did not signal READY (got: {ready!r})"
            )
        print("  [docling] Worker ready.", flush=True)

    def _drain_stderr(self):
        pass  # stderr discarded — no drain needed

    def convert(self, wsl_path: str, timeout: int = 180) -> dict:
        import threading
        self._ensure_started()
        self._proc.stdin.write(wsl_path + "\n")
        self._proc.stdin.flush()

        result = {}
        def _read():
            line = self._proc.stdout.readline()
            if line:
                result["data"] = json.loads(line)

        t = threading.Thread(target=_read, daemon=True)
        t.start()
        t.join(timeout=timeout)

        if t.is_alive():
            self._proc.kill()
            self._proc = None
            raise TimeoutError(f"Docling worker timed out after {timeout}s")

        # Check if worker died unexpectedly between files
        exit_code = self._proc.poll()
        if exit_code is not None:
            print(f"  [docling] Worker exited unexpectedly (code={exit_code})", flush=True)
            self._proc = None

        return result.get("data", {"ok": False, "error": "No output from worker"})

    def shutdown(self):
        if self._proc and self._proc.poll() is None:
            try:
                self._proc.stdin.close()
                self._proc.wait(timeout=5)
            except Exception:
                self._proc.kill()
        self._proc = None
        daemon_file = getattr(self, "_daemon_file", None)
        if daemon_file:
            try:
                os.unlink(daemon_file.name)
            except Exception:
                pass
            self._daemon_file = None


_worker = _DoclingWorker()

BATCH_SIZE = 40          # pages per docling batch for large PDFs
LARGE_PDF_THRESHOLD = 40 # PDFs with more pages than this are split


def _split_pdf_into_batches(file_path: str, batch_size: int = BATCH_SIZE):
    """
    Split a PDF into temp files of batch_size pages each.
    Returns list of (temp_path, start_page, end_page) tuples.
    Caller is responsible for deleting the temp files.
    """
    import pypdf
    reader = pypdf.PdfReader(file_path)
    total = len(reader.pages)
    batches = []
    for start in range(0, total, batch_size):
        end = min(start + batch_size, total)
        writer = pypdf.PdfWriter()
        for page_num in range(start, end):
            writer.add_page(reader.pages[page_num])
        tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        tmp.close()
        with open(tmp.name, "wb") as f:
            writer.write(f)
        batches.append((tmp.name, start + 1, end))
    return batches, total


def extract_with_docling(file_path: str, timeout: int = 180) -> Dict[str, Any]:
    """
    Send file to the persistent WSL docling worker.
    Large PDFs are split into BATCH_SIZE-page chunks and processed sequentially,
    with results merged into a single output.
    """
    import pypdf
    file_name = os.path.basename(file_path)

    # Count pages for PDFs to decide whether batching is needed
    total_pages = 0
    if file_path.lower().endswith(".pdf"):
        try:
            total_pages = len(pypdf.PdfReader(file_path).pages)
        except Exception:
            total_pages = 0

    use_batching = total_pages > LARGE_PDF_THRESHOLD

    if not use_batching:
        # Standard single-pass extraction
        try:
            wsl_path = _windows_path_to_wsl(os.path.abspath(file_path))
            print(f"  -> docling: {file_name} ({total_pages or '?'} pages)", flush=True)
            data = _worker.convert(wsl_path, timeout=timeout)
            if not data.get("ok"):
                raise RuntimeError(data.get("error", "unknown error"))
            chunks = data.get("chunks", [])
            print(f"  -> done: {len(chunks)} chunks", flush=True)
            return {"markdown": data["markdown"], "chunks": chunks}
        except Exception as e:
            print(f"  -> ERROR ({file_name}): {e}", flush=True)
            logger.error(f"Docling extraction failed for {file_path}: {e}")
            return {"error": str(e)}

    # ── Batched extraction for large PDFs ----------------------------------
    print(f"  -> docling: {file_name} ({total_pages} pages, batching every {BATCH_SIZE} pages)", flush=True)
    try:
        batches, total = _split_pdf_into_batches(file_path)
    except Exception as e:
        logger.error(f"PDF split failed for {file_path}: {e}")
        return {"error": str(e)}

    all_chunks = []
    all_markdown_parts = []
    batch_errors = []

    for batch_path, start_page, end_page in batches:
        try:
            wsl_path = _windows_path_to_wsl(os.path.abspath(batch_path))
            print(f"     batch pages {start_page}-{end_page} of {total}...", flush=True)
            data = _worker.convert(wsl_path, timeout=timeout)
            if not data.get("ok"):
                raise RuntimeError(data.get("error", "unknown error"))
            all_chunks.extend(data.get("chunks", []))
            all_markdown_parts.append(data.get("markdown", ""))
        except Exception as e:
            msg = f"batch pages {start_page}-{end_page}: {e}"
            print(f"     ERROR: {msg}", flush=True)
            logger.error(f"Docling batch error for {file_name} {msg}")
            batch_errors.append(msg)
        finally:
            try:
                os.unlink(batch_path)
            except Exception:
                pass

    if not all_chunks and not all_markdown_parts:
        return {"error": "; ".join(batch_errors) or "All batches failed"}

    combined_markdown = "\n\n".join(p for p in all_markdown_parts if p)
    print(f"  -> done: {len(all_chunks)} chunks across {len(batches)} batches", flush=True)
    if batch_errors:
        print(f"  -> {len(batch_errors)} batch(es) failed: {batch_errors}", flush=True)

    return {"markdown": combined_markdown, "chunks": all_chunks}


# ── pandas (XLSX / XLS) ----------------------------------------------------

def extract_xlsx(file_path: str) -> Dict[str, Any]:
    """Load all sheets, normalise NaNs to empty strings."""
    try:
        dfs = pd.read_excel(file_path, sheet_name=None)
        return {
            sheet: df.fillna("").astype(str).to_dict(orient="records")
            for sheet, df in dfs.items()
        }
    except Exception as e:
        logger.error(f"XLSX extraction failed for {file_path}: {e}")
        return {"error": str(e)}


# ── python-pptx (PPTX / PPT) -----------------------------------------------

def extract_pptx(file_path: str) -> Dict[str, Any]:
    """Extract slide text and tables from a PowerPoint presentation."""
    try:
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            tables = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
                if shape.has_table:
                    tables.append([
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ])
            slides.append({"slide_number": i, "texts": texts, "tables": tables})
        return {"slides": slides}
    except Exception as e:
        logger.error(f"PPTX extraction failed for {file_path}: {e}")
        return {"error": str(e)}


# ── Raw text (XML, TXT, fallback) ------------------------------------------

def extract_text(file_path: str) -> Dict[str, Any]:
    """Read file as plain text."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return {"raw_text": f.read()}
    except Exception as e:
        logger.error(f"Text extraction failed for {file_path}: {e}")
        return {"error": str(e)}


# ── Router -----------------------------------------------------------------

def process_single_document(file_path: str, timeout: int = 180, is_retry: bool = False) -> bool:
    """Route file to correct extractor, save JSON to extracted_docs/. Returns True if processed."""
    ext = _get_ext(file_path)

    if ext in SKIP_EXTENSIONS:
        return False

    folder_path = os.path.dirname(file_path)
    file_name = os.path.basename(file_path)
    extracted_folder = os.path.join(folder_path, "extracted_docs")
    os.makedirs(extracted_folder, exist_ok=True)

    output_json_path = os.path.join(extracted_folder, f"{file_name}.json")
    failed_marker_path = os.path.join(extracted_folder, f"{file_name}.failed")

    if os.path.exists(output_json_path):
        logger.debug(f"Already extracted, skipping: {output_json_path}")
        return False

    MAX_ATTEMPTS = 2  # permanently skip after this many failures

    if os.path.exists(failed_marker_path):
        try:
            marker = json.loads(open(failed_marker_path, encoding="utf-8").read())
            attempts = marker.get("attempts", 1)
        except Exception:
            attempts = 1

        if attempts >= MAX_ATTEMPTS:
            logger.debug(f"Permanently skipping after {attempts} failed attempts: {file_name}")
            return False

        if not is_retry:
            logger.debug(f"Previously failed, skipping (attempt {attempts}/{MAX_ATTEMPTS}): {file_name}")
            return False

    logger.info(f"Extracting {file_name} ...")

    if ext in DOCLING_FORMATS:
        content = extract_with_docling(file_path, timeout=timeout)
    elif ext in XLSX_FORMATS:
        content = extract_xlsx(file_path)
    elif ext in PPTX_FORMATS:
        content = extract_pptx(file_path)
    else:
        content = extract_text(file_path)

    if content is None:
        return False

    # Write a .failed marker so this file is skipped on future runs
    if list(content.keys()) == ["error"]:
        error_msg = content["error"]
        # Read existing attempt count if marker already exists
        prior_attempts = 0
        if os.path.exists(failed_marker_path):
            try:
                prior = json.loads(open(failed_marker_path, encoding="utf-8").read())
                prior_attempts = prior.get("attempts", 1)
            except Exception:
                prior_attempts = 1
        new_attempts = prior_attempts + 1
        logger.warning(f"Extraction failed for {file_name} (attempt {new_attempts}): {error_msg}")
        try:
            with open(failed_marker_path, "w", encoding="utf-8") as f:
                json.dump({"error": error_msg, "attempts": new_attempts}, f)
        except Exception:
            pass
        return False

    output_data = {
        "source_file": file_name,
        "extension": ext.lstrip("."),
        "extraction_timestamp": datetime.utcnow().isoformat() + "Z",
        "content": content,
    }

    try:
        with open(output_json_path, "w", encoding="utf-8") as f:
            json.dump(output_data, f, ensure_ascii=False, indent=2, cls=SafeEncoder)
        return True
    except Exception as e:
        logger.error(f"Failed to save JSON for {file_name}: {e}")
        return False


# ── Runner -----------------------------------------------------------------

def run_document_extraction() -> dict:
    """
    Walk all tender document folders and extract text from each file.

    NOTE: This is the legacy all-in-one runner. The preferred workflow is:
      1. python pre_process_workflow/01_classify.py
      2. python pre_process_workflow/02_build_queue.py
      3. python pre_process_workflow/03_extract.py
    """
    docs_dir = os.path.join(config.TENDERS_OUTPUT_DIRECTORY, "documents")
    if not os.path.exists(docs_dir):
        logger.warning(f"Documents directory not found at {docs_dir}")
        return {"total_files_scanned": 0, "newly_processed": 0, "skipped": 0, "errors": 0}

    logger.info(f"Scanning {docs_dir} for documents ...")

    all_files = []
    for root, dirs, files in os.walk(docs_dir):
        dirs[:] = [d for d in dirs if d != "extracted_docs"]
        for file in files:
            all_files.append(os.path.join(root, file))

    total = len(all_files)
    processed = 0
    skipped = 0
    errors = 0

    print(f"Found {total} files to scan.", flush=True)

    for i, file_path in enumerate(all_files, start=1):
        file_name = os.path.basename(file_path)
        print(f"[{i}/{total}] {file_name}", flush=True)
        try:
            if process_single_document(file_path):
                processed += 1
            else:
                skipped += 1
        except Exception as e:
            print(f"  -> ERROR: {e}", flush=True)
            logger.error(f"Failed to process {file_path}: {e}")
            errors += 1

    # ── Retry pass: re-attempt any .failed files with a longer timeout ----------
    failed_markers = []
    for root, dirs, files in os.walk(docs_dir):
        for file in files:
            if file.endswith(".failed"):
                failed_markers.append(os.path.join(root, file))

    if failed_markers:
        print(f"\nRetrying {len(failed_markers)} previously failed file(s) with extended timeout (300s)...", flush=True)
        retried = 0
        retry_errors = 0
        for i, marker_path in enumerate(failed_markers, start=1):
            # Marker lives in extracted_docs/filename.pdf.failed
            # Original file is one level up: ../filename.pdf
            marker_name = os.path.basename(marker_path)
            original_name = marker_name[: -len(".failed")]
            extracted_dir = os.path.dirname(marker_path)
            original_path = os.path.join(os.path.dirname(extracted_dir), original_name)
            file_name = original_name
            print(f"  [retry {i}/{len(failed_markers)}] {file_name}", flush=True)
            if not os.path.exists(original_path):
                print(f"    -> original file missing, removing marker", flush=True)
                os.unlink(marker_path)
                continue
            try:
                if process_single_document(original_path, timeout=300, is_retry=True):
                    os.unlink(marker_path)  # clear marker on success
                    retried += 1
                    processed += 1
                else:
                    retry_errors += 1
            except Exception as e:
                print(f"    -> ERROR: {e}", flush=True)
                logger.error(f"Retry failed for {original_path}: {e}")
                retry_errors += 1

        errors += retry_errors
        print(f"Retry complete: {retried} recovered, {retry_errors} still failing.", flush=True)

    result = {
        "total_files_scanned": total,
        "newly_processed": processed,
        "skipped": skipped,
        "errors": errors,
    }
    _worker.shutdown()
    logger.info(f"Extraction complete: {result}")
    return result


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    run_document_extraction()
