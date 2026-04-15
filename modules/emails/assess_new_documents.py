"""
VLM-based assessment of newly downloaded email-update documents.

Workflow:
  1. Quick-extract content from each new/changed file:
       PDF          -> first N pages rendered as base64 images (fitz/PyMuPDF)
       DOCX         -> plain text (python-docx)
       XLSX / XLS   -> sheet text (pandas)
       PPTX / PPT   -> slide text (python-pptx)
       TXT / XML    -> raw read
  2. Fetch the current analysis summary from gojep_analysis_results
  3. Send both to qwen3_vl_32b via OpenRouter
  4. VLM returns { needs_reanalysis: bool, reason: str }

If extraction fails for a file, or the VLM call fails, the function
returns needs_reanalysis=True as a safe default.
"""

from __future__ import annotations

import base64
import json
import logging
import os
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# Pages per PDF to send to the VLM — enough to capture key content without
# blowing the context window on large documents
MAX_PDF_PAGES = 4

# Max characters of extracted text per file
MAX_TEXT_CHARS = 6_000

_ASSESSMENT_SYSTEM_PROMPT = """\
You are a procurement analyst assessing whether newly received documents for a \
Government of Jamaica tender contain information that materially changes the \
existing analysis.

You will be given:
  1. The current analysis of the tender (JSON)
  2. Content extracted from newly received documents

Determine whether the new documents introduce any material changes that would \
require a full re-analysis. Material changes include:
  - New or changed scope of work or technical requirements
  - New or changed eligibility or experience requirements
  - New or changed evaluation criteria or weightings
  - New or changed mandatory submission documents
  - Significant amendments to contract terms or conditions
  - New lots or contract packages

The following do NOT require re-analysis:
  - Deadline reminders or confirmations
  - Venue or logistical clarifications with no scope impact
  - Minor typographical corrections
  - Acknowledgement notices

Return ONLY a valid JSON object:
{
  "needs_reanalysis": true | false,
  "reason": "One sentence explaining the decision"
}
"""


# ── File content extraction ───────────────────────────────────────────────────

def _extract_pdf_pages_as_images(path: Path) -> list[dict]:
    """Render up to MAX_PDF_PAGES pages as base64 PNG images for VLM input."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF not installed — cannot render PDF pages for %s", path.name)
        return []

    images = []
    try:
        doc = fitz.open(str(path))
        for i in range(min(MAX_PDF_PAGES, len(doc))):
            page = doc[i]
            pix = page.get_pixmap(dpi=120)
            png_bytes = pix.tobytes("png")
            b64 = base64.b64encode(png_bytes).decode("utf-8")
            images.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{b64}"},
            })
        doc.close()
    except Exception as e:
        logger.warning("PDF render failed for %s: %s", path.name, e)
    return images


def _extract_docx_text(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return text[:MAX_TEXT_CHARS]
    except Exception as e:
        logger.warning("DOCX extraction failed for %s: %s", path.name, e)
        return ""


def _extract_xlsx_text(path: Path) -> str:
    try:
        import pandas as pd
        xf = pd.ExcelFile(str(path))
        parts = []
        for sheet in xf.sheet_names:
            df = pd.read_excel(xf, sheet_name=sheet, header=None)
            parts.append(f"[Sheet: {sheet}]\n{df.to_string(index=False, header=False)}")
        return "\n\n".join(parts)[:MAX_TEXT_CHARS]
    except Exception as e:
        logger.warning("XLSX extraction failed for %s: %s", path.name, e)
        return ""


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = [
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                parts.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
        return "\n\n".join(parts)[:MAX_TEXT_CHARS]
    except Exception as e:
        logger.warning("PPTX extraction failed for %s: %s", path.name, e)
        return ""


def _extract_file_content(path: Path) -> tuple[str, list[dict]]:
    """
    Extract content from a file.
    Returns (text_content, image_parts) — PDFs return images, others return text.
    """
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        images = _extract_pdf_pages_as_images(path)
        text = f"[PDF: {path.name} — {len(images)} page(s) rendered as images]" if images else ""
        return text, images

    if suffix in (".docx", ".doc"):
        return _extract_docx_text(path), []

    if suffix in (".xlsx", ".xls"):
        return _extract_xlsx_text(path), []

    if suffix in (".pptx", ".ppt"):
        return _extract_pptx_text(path), []

    # TXT, XML, and anything else — read as raw text
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return text[:MAX_TEXT_CHARS], []
    except Exception as e:
        logger.warning("Text read failed for %s: %s", path.name, e)
        return "", []


# ── Current analysis fetch ────────────────────────────────────────────────────

def _fetch_current_analysis(resource_id: str) -> dict[str, Any]:
    """Fetch the existing analysis from gojep_analysis_results."""
    try:
        from config import settings as config
        from db.client.supabase_client import SupabaseClient
        db = SupabaseClient()
        rows = (
            db.supabase.table(config.SUPABASE_TABLE_ANALYSIS_RESULTS)
            .select(
                "contract_title,procuring_entity,contract_type,scope_of_work,"
                "eligibility_requirements,experience_requirements,"
                "financial_requirements,mandatory_documents,"
                "evaluation_criteria,key_milestones,special_conditions"
            )
            .eq("resource_id", resource_id)
            .limit(1)
            .execute()
            .data or []
        )
        return rows[0] if rows else {}
    except Exception as e:
        logger.warning("Could not fetch current analysis for %s: %s", resource_id, e)
        return {}


# ── VLM assessment ────────────────────────────────────────────────────────────

def assess_new_documents_for_reanalysis(
    new_file_paths: list[str],
    resource_id: str,
    update_type: str = "unknown",
) -> dict[str, Any]:
    """
    Assess whether newly downloaded documents warrant a full re-analysis.

    Args:
        new_file_paths:  Absolute paths to the new/changed files from hash comparison.
        resource_id:     Tender resource_id — used to fetch the current analysis.
        update_type:     Email update type for logging context.

    Returns:
        {
            "needs_reanalysis": bool,
            "reason": str,
            "assessed_files": [str],
        }
    """
    from config import settings as config
    import requests as _requests

    if not new_file_paths:
        return {"needs_reanalysis": False, "reason": "no new files to assess", "assessed_files": []}

    # ── Build message content from new files ──────────────────────────────────
    file_sections: list[str] = []
    all_image_parts: list[dict] = []
    assessed_files: list[str] = []

    for fpath in new_file_paths:
        path = Path(fpath)
        if not path.exists():
            continue
        text, images = _extract_file_content(path)
        if text or images:
            file_sections.append(f"--- {path.name} ---\n{text}" if text else f"--- {path.name} ---")
            all_image_parts.extend(images)
            assessed_files.append(path.name)
        else:
            logger.warning("  No content extracted from %s — skipping", path.name)

    if not assessed_files:
        logger.warning("  No extractable content from new files — defaulting to needs_reanalysis=True")
        return {
            "needs_reanalysis": True,
            "reason": "could not extract content from new files — re-analysis triggered as safe default",
            "assessed_files": [],
        }

    # ── Fetch current analysis ────────────────────────────────────────────────
    current_analysis = _fetch_current_analysis(resource_id)
    if not current_analysis:
        logger.warning("  No existing analysis found for %s — defaulting to needs_reanalysis=True", resource_id)
        return {
            "needs_reanalysis": True,
            "reason": "no existing analysis to compare against — re-analysis triggered",
            "assessed_files": assessed_files,
        }

    # ── Build VLM prompt ──────────────────────────────────────────────────────
    current_analysis_text = json.dumps(current_analysis, indent=2, ensure_ascii=False)
    new_docs_text = "\n\n".join(file_sections)

    user_text = (
        f"UPDATE TYPE: {update_type}\n\n"
        f"CURRENT ANALYSIS:\n{current_analysis_text}\n\n"
        f"NEW/UPDATED DOCUMENTS ({len(assessed_files)} file(s)):\n{new_docs_text}"
    )

    # Build message content — text first, then images (VLM interleaved input)
    user_content: list[dict] = [{"type": "text", "text": user_text}]
    user_content.extend(all_image_parts)

    payload = {
        "model":       config.OPENROUTER_MODELS[config.ANALYSIS_MODEL],
        "messages":    [
            {"role": "system", "content": _ASSESSMENT_SYSTEM_PROMPT},
            {"role": "user",   "content": user_content},
        ],
        "max_tokens":  256,
        "temperature": 0.0,
    }

    # ── Call VLM ──────────────────────────────────────────────────────────────
    try:
        resp = _requests.post(
            config.OPENROUTER_URL,
            headers=config.OPENROUTER_HEADERS,
            json=payload,
            timeout=120,
        )
        resp.raise_for_status()
        content = resp.json()["choices"][0]["message"]["content"].strip()

        # Strip markdown fences if present
        if content.startswith("```"):
            import re
            content = re.sub(r"^```[^\n]*\n?", "", content)
            content = re.sub(r"\n?```$", "", content)

        result = json.loads(content)
        needs = bool(result.get("needs_reanalysis", True))
        reason = result.get("reason", "")

        logger.info(
            "  VLM assessment for %s (%s): needs_reanalysis=%s — %s",
            resource_id, update_type, needs, reason,
        )
        return {
            "needs_reanalysis": needs,
            "reason":           reason,
            "assessed_files":   assessed_files,
        }

    except Exception as e:
        logger.warning(
            "  VLM assessment failed for %s: %s — defaulting to needs_reanalysis=True", resource_id, e
        )
        return {
            "needs_reanalysis": True,
            "reason":           f"VLM assessment failed ({e}) — re-analysis triggered as safe default",
            "assessed_files":   assessed_files,
        }
